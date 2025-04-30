import os
from typing import Any, Dict, List, Optional

from app.filebase.parsers.base_parser import BaseParser
from app.filebase.vector.file_chunk import FileChunk
from app.logger import get_logger

logger = get_logger(__name__)


class PowerPointParser(BaseParser):
    """
    PowerPoint文件解析器，处理.pptx和.ppt文件
    """

    # 支持的文件扩展名
    SUPPORTED_EXTENSIONS: List[str] = ['.pptx', '.ppt']

    def can_handle(self, file_path: str) -> bool:
        """
        检查是否可以处理该文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            bool: 是否可以处理
        """
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.SUPPORTED_EXTENSIONS

    def parse(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """
        解析PowerPoint文件
        
        Args:
            file_path: 文件路径
            metadata: 附加元数据
            
        Returns:
            Dict[str, Any]: 解析结果，包括元数据、内容和FileChunk对象列表
        """
        if metadata is None:
            metadata = {}

        # 添加解析器信息到元数据
        metadata['parser'] = 'powerpoint'
        metadata['file_type'] = 'powerpoint'

        try:
            # 使用python-pptx库解析PPT
            try:
                from pptx import Presentation
                prs = Presentation(file_path)

                # 获取幻灯片数量
                slide_count = len(prs.slides)
                metadata['slide_count'] = slide_count

                # 提取每张幻灯片的文本
                all_content = []
                file_chunks = []

                for i, slide in enumerate(prs.slides):
                    slide_content = [f"--- 幻灯片 {i+1} ---"]

                    # 提取标题
                    if slide.shapes.title and slide.shapes.title.text:
                        slide_content.append(f"标题: {slide.shapes.title.text}")

                    # 提取所有文本框中的文本
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_texts.append(shape.text)

                    if slide_texts:
                        slide_content.append("文本内容:")
                        slide_content.extend(slide_texts)
                    else:
                        slide_content.append("[无文本内容]")

                    # 合并当前幻灯片的内容
                    current_slide_content = "\n".join(slide_content)
                    all_content.append(current_slide_content)

                    # 为每张幻灯片创建一个 FileChunk
                    chunk_metadata = {
                        'slide_number': i + 1,
                        'content_type': 'ppt_slide'
                    }

                    file_chunk = FileChunk(
                        text=current_slide_content,
                        file_metadata=metadata,
                        chunk_metadata=chunk_metadata,
                        chunk_index=i,
                        total_chunks=slide_count
                    )
                    file_chunks.append(file_chunk)

                # 合并所有幻灯片内容
                content = "\n\n".join(all_content)

                return {
                    'metadata': metadata,
                    'content': content,
                    'chunks': file_chunks
                }
            except ImportError:
                logger.warning("未安装python-pptx库，无法解析PPT文件内容")
                return {
                    'metadata': metadata,
                    'content': "错误：解析PPT文件需要安装python-pptx库（pip install python-pptx）",
                    'chunks': []
                }

        except Exception as e:
            logger.error(f"解析PowerPoint文件 {file_path} 时出错: {e!s}")
            return {
                'metadata': metadata,
                'content': f"PowerPoint content from {file_path} 解析失败: {e!s}",
                'chunks': []
            } 
